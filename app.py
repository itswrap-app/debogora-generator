import streamlit as st
import pandas as pd
from datetime import datetime, date, timedelta
import io
import os
import subprocess
import base64
import time
import re
import glob
import requests
import json
import shutil
import uuid

# Biblioteki Google
from google.oauth2.service_account import Credentials as SACredentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload, MediaIoBaseUpload
from pptx import Presentation
from pypdf import PdfWriter, PdfReader

# ReportLab
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image as RLImage
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import reportlab.rl_config

reportlab.rl_config.warnOnMissingFontGlyphs = 0

st.set_page_config(page_title="PLANER OFERT - Dwór Dębogóra", layout="wide")

CI = {
    "dark_green": "#00622f",
    "light_green": "#e8ece6",
    "gray": "#333333",
    "white": "#ffffff"
}

# --- IDENTYFIKATORY GOOGLE ---
ROOT_FOLDER_ID = "1tU6mo1YWpTep8vl5CRR5DhsZAINeWnHz"  
BAZA_OFERT_FOLDER_ID = "1i_a2UkK73ixyvMBe5l9SkE5vpqAu6he5" 
BAZA_OFERT_SHEET_ID = "TUTAJ_WKLEJ_ID_ARKUSZA" 

# =========================================================
# MAPOWANIE POKOI HOTRES
# =========================================================
HOTRES_ROOM_MAP = {
    29952: "Krovacja - cały kompleks",
    37951: "Muuu 6",
    37950: "Muuu 5",
    25074: "Muuu 4",
    25073: "Muuu 3",
    25072: "Muuu 2",
    25071: "Muuu 1",
    27589: "Ognisko #1",
    27588: "Strefa relaksu - balia #1",
    31294: "Strefa relaksu - balia #2",
    31230: "Łaźnia eventowa",
    23704: "Oranżeria",
    23698: "Sala kominkowa",
    31228: "Sala wielofunkcyjna",
    23705: "Sala bankietowa",
    22144: "Pokój nr 2",
    22146: "Pokój nr 3",
    22147: "Pokój nr 4",
    22148: "Pokój nr 5",
    22149: "Pokój nr 6",
    22150: "Pokój nr 7",
    22151: "Pokój nr 8",
    22152: "Pokój nr 9",
    22153: "Pokój nr 10",
    22154: "Pokój nr 11",
    22156: "Pokój nr 12",
    31229: "Dwór - cały obiekt"
}

HOTRES_ROOM_NAME_TO_ID = {v: k for k, v in HOTRES_ROOM_MAP.items()}

# --- BEZPIECZNA INICJALIZACJA STANU APLIKACJI ---
if "klient_imie" not in st.session_state: st.session_state.klient_imie = ""
if "firma_n" not in st.session_state: st.session_state.firma_n = ""
if "nip_n" not in st.session_state: st.session_state.nip_n = ""
if "telefon_n" not in st.session_state: st.session_state.telefon_n = ""
if "email_n" not in st.session_state: st.session_state.email_n = ""
if "loaded_pozycje" not in st.session_state: st.session_state.loaded_pozycje = None
if "agenda_custom_text" not in st.session_state: st.session_state.agenda_custom_text = ""
if "l_osob_total" not in st.session_state: st.session_state.l_osob_total = 10
if "szczegoly_zajetosci" not in st.session_state: st.session_state.szczegoly_zajetosci = {}
if "wybrane_p" not in st.session_state: st.session_state.wybrane_p = []
if "wybrane_d" not in st.session_state: st.session_state.wybrane_d = []
if "pdf_page_list" not in st.session_state: st.session_state.pdf_page_list = []

# Nowe klucze sesji do odtwarzania UI z Archiwum
if "wyz_sel_key" not in st.session_state: st.session_state.wyz_sel_key = []
if "spa_sel_key" not in st.session_state: st.session_state.spa_sel_key = []
if "atr_sel_key" not in st.session_state: st.session_state.atr_sel_key = []
if "biz_sel_key" not in st.session_state: st.session_state.biz_sel_key = []

# --- INTELIGENTNE ŁADOWANIE CZCIONEK Z DYSKU ---
FONT_HEADER = 'Helvetica-Bold'
FONT_TEXT = 'Helvetica'
FONT_TEXT_BOLD = 'Helvetica-Bold'

def install_fonts_for_libreoffice():
    try:
        fonts_dir = os.path.expanduser('~/.fonts')
        os.makedirs(fonts_dir, exist_ok=True)
        copied = False
        for font_file in glob.glob('*.ttf'):
            target_path = os.path.join(fonts_dir, font_file)
            if not os.path.exists(target_path):
                shutil.copy(font_file, target_path)
                copied = True
        if copied:
            subprocess.run(["fc-cache", "-f"], capture_output=True)
    except Exception:
        pass

def register_custom_fonts():
    global FONT_HEADER, FONT_TEXT, FONT_TEXT_BOLD
    try:
        if os.path.exists('Lora-Bold.ttf'):
            pdfmetrics.registerFont(TTFont('Lora-Bold', 'Lora-Bold.ttf'))
            FONT_HEADER = 'Lora-Bold'
    except Exception:
        pass
    
    try:
        if os.path.exists('PTSans-Regular.ttf'):
            pdfmetrics.registerFont(TTFont('PTSans-Regular', 'PTSans-Regular.ttf'))
            FONT_TEXT = 'PTSans-Regular'
    except Exception:
        pass
        
    try:
        if os.path.exists('PTSans-Bold.ttf'):
            pdfmetrics.registerFont(TTFont('PTSans-Bold', 'PTSans-Bold.ttf'))
            FONT_TEXT_BOLD = 'PTSans-Bold'
    except Exception:
        pass

register_custom_fonts()

st.markdown(f"""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Lora:wght@400;700&family=PT+Sans:wght@400;700&display=swap');
    .stApp {{ background-color: {CI['white']}; font-family: 'PT Sans', sans-serif; }}
    h1, h2, h3, h4 {{ font-family: 'Lora', serif !important; color: {CI['dark_green']} !important; font-weight: 700 !important; }}
    div[data-testid="stVerticalBlock"] > div > div[data-testid="stVerticalBlock"] {{
        background-color: {CI['light_green']}; padding: 2rem; border-left: 5px solid {CI['dark_green']}; margin-bottom: 1.5rem;
    }}
    div.stButton > button {{
        background-color: {CI['dark_green']} !important; color: white !important;
        border-radius: 0px !important; font-family: 'Lora', serif !important; padding: 0.8rem 3rem !important;
        text-transform: uppercase; letter-spacing: 2px;
    }}
    div.stButton > button:hover {{ background-color: {CI['gray']} !important; }}
    </style>
""", unsafe_allow_html=True)

# --- GOOGLE DRIVE & SHEETS LOGIC ---
def get_google_credentials():
    info = st.secrets["gcp_service_account"]
    return SACredentials.from_service_account_info(
        info,
        scopes=['https://www.googleapis.com/auth/drive', 'https://www.googleapis.com/auth/spreadsheets']
    )

@st.cache_resource
def get_drive_service():
    return build('drive', 'v3', credentials=get_google_credentials())

@st.cache_resource
def get_sheets_service():
    return build('sheets', 'v4', credentials=get_google_credentials())

def save_offer_to_sheet(sheet_id, meta_data):
    try:
        service = get_sheets_service()
        row = [
            datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            meta_data.get("klient_imie", ""),
            meta_data.get("firma_n", ""),
            meta_data.get("nip_n", ""),
            meta_data.get("telefon_n", ""),
            meta_data.get("email_n", ""),
            meta_data.get("marka_oferty", ""),
            meta_data.get("typ_klienta", ""),
            meta_data.get("l_osob_total", ""),
            meta_data.get("final_agenda_text", ""),
            json.dumps(meta_data.get("pozycje", []), ensure_ascii=False)
        ]
        body = {'values': [row]}
        service.spreadsheets().values().append(
            spreadsheetId=sheet_id, 
            range="Arkusz1!A:K", 
            valueInputOption="USER_ENTERED", 
            body=body
        ).execute()
    except Exception as e:
        st.warning(f"Nie udało się dopisać rekordu do arkusza: {e}")

def get_offers_from_sheet(sheet_id):
    try:
        service = get_sheets_service()
        result = service.spreadsheets().values().get(
            spreadsheetId=sheet_id, 
            range="Arkusz1!A:K"
        ).execute()
        return result.get('values', [])
    except Exception:
        return []

def upload_file_to_drive(file_bytes, filename, folder_id, mimetype='application/pdf'):
    try:
        service = get_drive_service()
        media = MediaIoBaseUpload(io.BytesIO(file_bytes), mimetype=mimetype, resumable=True)
        file_metadata = {'name': filename, 'parents': [folder_id]}
        file = service.files().create(body=file_metadata, media_body=media, fields='id, webViewLink', supportsAllDrives=True).execute()
        try:
            service.permissions().create(
                fileId=file.get('id'), 
                body={'type': 'anyone', 'role': 'reader'}, 
                supportsAllDrives=True
            ).execute()
        except Exception: 
            pass
        return file
    except Exception:
        return None

@st.cache_data(ttl=600, show_spinner=False)
def fetch_all_debogora_files(root_id):
    try:
        service = get_drive_service()
        all_files = []
        folders_to_search = [root_id]
        
        while folders_to_search:
            current_folder = folders_to_search.pop(0)
            query = f"'{current_folder}' in parents and trashed = false"
            for _ in range(3):
                try:
                    request = service.files().list(
                        q=query, 
                        fields="nextPageToken, files(id, name, mimeType, webViewLink)", 
                        pageSize=1000, 
                        supportsAllDrives=True, 
                        includeItemsFromAllDrives=True
                    )
                    while request is not None:
                        results = request.execute()
                        files = results.get('files', [])
                        for f in files:
                            if f['mimeType'] == 'application/vnd.google-apps.folder':
                                folders_to_search.append(f['id'])
                            else:
                                all_files.append(f)
                        request = service.files().list_next(request, results)
                    break
                except Exception:
                    time.sleep(1.0)
        return all_files
    except Exception:
        return []

def download_file(file_id, retries=3):
    service = get_drive_service()
    for attempt in range(retries):
        try:
            request = service.files().get_media(fileId=file_id, supportsAllDrives=True)
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request, chunksize=256*1024)
            done = False
            while done is False:
                status, done = downloader.next_chunk()
            fh.seek(0)
            return fh
        except Exception as e:
            if attempt == retries - 1: 
                raise e
            time.sleep(1.0)

def update_file_on_drive(file_id, df, file_name):
    service = get_drive_service()
    buffer = io.BytesIO()
    if 'xlsx' in file_name.lower():
        df.to_excel(buffer, index=False, engine='openpyxl')
        mimetype = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    else:
        df.to_csv(buffer, index=False, encoding='utf-8')
        mimetype = 'text/csv'
    buffer.seek(0)
    media = MediaIoBaseUpload(buffer, mimetype=mimetype, resumable=True)
    service.files().update(fileId=file_id, media_body=media, supportsAllDrives=True).execute()

# --- ZAMIANA TEKSTU W PPTX ---
def process_shape(shape, replacements):
    if hasattr(shape, "text_frame") and shape.text_frame is not None:
        for paragraph in shape.text_frame.paragraphs:
            if not paragraph.runs: 
                continue
            full_text = "".join(run.text for run in paragraph.runs)
            replaced = False
            for k, v in replacements.items():
                if k in full_text:
                    full_text = full_text.replace(k, str(v))
                    replaced = True
            if replaced and paragraph.runs:
                paragraph.runs[0].text = full_text
                for i in range(1, len(paragraph.runs)):
                    paragraph.runs[i].text = ""

    if hasattr(shape, "shapes"):
        for subshape in shape.shapes: 
            process_shape(subshape, replacements)
            
    if hasattr(shape, "has_table") and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells: 
                process_shape(cell, replacements)

def replace_text_in_pptx(prs, replacements):
    for slide in prs.slides:
        for shape in slide.shapes: 
            process_shape(shape, replacements)

# --- MAPOWANIE NAZW I WYSZUKIWANIE PLIKÓW ---
def normalize_pl(text):
    rep = {'ą':'a', 'ć':'c', 'ę':'e', 'ł':'l', 'ń':'n', 'ó':'o', 'ś':'s', 'ź':'z', 'ż':'z'}
    res = str(text).lower()
    for k, v in rep.items(): 
        res = res.replace(k, v)
    return re.sub(r'[\W_]+', '', res)

SYNONYMS = {
    "okładka": "okładka", 
    "powitalna": "karta powitalna", 
    "zakwaterowanie_dwor": "dwór 3", 
    "zakwaterowanie_domki": "zakwaterowanie domki",
    "zakwaterowanie_oba": "dwór i domki", 
    "wyżywienie": "wyżywienie", 
    "serwis kawowy": "serwis kawowy", 
    "wiejskie jadło": "wiejskie jadło",
    "rozszerzonym menu": "kolacja z rozszerzonym menu", 
    "atrakcje_wstęp": "atrakcje", 
    "wycena": "wycena", 
    "agenda": "agenda", 
    "kontakt": "AsystentAI_kontakt", 
    "uklad_debogora": "AsystentAI_układ łóżek_Dębogóra", 
    "uklad_krovacja": "AsystentAI_układ łóżek_Krovacja",
    "Złodziej Krów": "złodziej", 
    "Skarby": "skarby", 
    "Safari_Standard": "safari dla grup_standard", 
    "Safari_Rozszerzona": "safari dla grup_rozszerzona", 
    "Seans saunowy": "saunowy", 
    "Sauna olchowa": "olchowa", 
    "Staw": "staw", 
    "Balia": "balia", 
    "Sauny": "sauny", 
    "Masaże": "masaż", 
    "Paintball": "paintball", 
    "Spływ kajakowy": "kajak",
    "Rowery": "rowery", 
    "Ognisko": "ognisko", 
    "Punkt widokowy": "widokowy", 
    "Łączka cielaczków": "cielacz",
    "Atrakcje na wodzie": "wodzie", 
    "Złów i wypuść": "złów", 
    "Grzybobranie": "grzyb", 
    "Roztańczony las": "roztańczony",
    "Drawieński PN": "drawieński", 
    "Blok konferencyjny": "konferencyjny", 
    "Wynajem sali": "sali", 
    "Przejazd": "przejazd"
}

def get_file_by_keyword(keyword, all_files):
    norm_search = normalize_pl(SYNONYMS.get(keyword, keyword))
    matches = [f for f in all_files if norm_search in normalize_pl(f['name'])]
    
    if keyword == "atrakcje_wstęp":
        matches = [f for f in matches if "wodzi" not in normalize_pl(f['name'])]
        
    if matches:
        matches.sort(key=lambda x: len(x['name']))
        prev_matches = [f for f in matches if 'prev' in f['name'].lower()]
        if prev_matches:
            return prev_matches[0]
        else:
            return matches[0]
    return None

def add_file_to_merger(merger, keyword, all_files, open_streams, missing_cards, added_file_ids, session_uid, replacements=None):
    if not keyword: 
        return
    file_obj = get_file_by_keyword(keyword, all_files)
    if file_obj:
        if file_obj['id'] in added_file_ids: 
            return
        try:
            fh = download_file(file_obj['id'])
            fname = file_obj['name'].lower()
            if 'ppt' in fname or 'presentation' in file_obj['mimeType']:
                temp_ppt = f"temp_{session_uid}_{file_obj['id']}.pptx"
                temp_pdf = f"temp_{session_uid}_{file_obj['id']}.pdf"
                lo_profile_dir = f"./lo_profile_{session_uid}"
                lo_profile_flag = f"-env:UserInstallation=file://{os.path.abspath(lo_profile_dir)}"
                
                with open(temp_ppt, "wb") as f: 
                    f.write(fh.getvalue())
                
                if replacements:
                    prs = Presentation(temp_ppt)
                    replace_text_in_pptx(prs, replacements)
                    prs.save(temp_ppt)
                    
                res = subprocess.run(["libreoffice", lo_profile_flag, "--headless", "--convert-to", "pdf", temp_ppt], capture_output=True, text=True)
                if res.returncode != 0: 
                    raise Exception(f"LibreOffice błąd: {res.stderr}")
                
                with open(temp_pdf, "rb") as f: 
                    pdf_bytes = f.read()
                pdf_stream = io.BytesIO(pdf_bytes)
                
                try: 
                    os.remove(temp_ppt)
                except: 
                    pass
                try: 
                    os.remove(temp_pdf)
                except: 
                    pass
                
            else:
                pdf_stream = fh
                
            open_streams.append(pdf_stream)
            merger.append(PdfReader(pdf_stream, strict=False))
            added_file_ids.add(file_obj['id'])
        except Exception as e:
            st.error(f"⚠️ Pominięto '{keyword}'. Błąd pliku '{file_obj['name']}': {e}")
            missing_cards.append(keyword)
    else:
        if keyword not in missing_cards: 
            missing_cards.append(keyword)

def safe_str(text): 
    if pd.isna(text):
        return ""
    else:
        return str(text).strip()

def get_price_data(usluga_name, df):
    search_names = [usluga_name.lower()]
    if "złodziej" in usluga_name.lower() or "łowcy" in usluga_name.lower(): 
        search_names.extend(["złodziej krów", "łowcy krów"])
        
    if df is None or df.empty: 
        return {"cena": 0, "czas": 0.0, "min_start": 9.0, "max_start": 18.0}
        
    try:
        col_name = next((c for c in df.columns if 'nazwa' in c.lower() or 'usługa' in c.lower() or 'usluga' in c.lower()), None)
        col_price = next((c for c in df.columns if 'cena' in c.lower()), None)
        col_czas = next((c for c in df.columns if 'długość' in c.lower() or 'dlugosc' in c.lower()), None)
        col_kiedy = next((c for c in df.columns if 'kiedy' in c.lower() or 'zacząć' in c.lower() or 'zaczac' in c.lower()), None)
        
        if col_name and col_price:
            match = df[df[col_name].astype(str).str.strip().str.lower().isin(search_names)]
            if match.empty: 
                match = df[df[col_name].astype(str).str.lower().str.contains(usluga_name.lower()[:5], na=False)]
                
            if not match.empty:
                val = match.iloc[0][col_price]
                if pd.notna(val):
                    cena = float(val)
                else:
                    cena = 0.0
                    
                czas_num = 1.0
                if col_czas and pd.notna(match.iloc[0][col_czas]):
                    czas_str = str(match.iloc[0][col_czas]).lower().replace('h', '').replace('godz.', '').replace('godz', '').strip()
                    if czas_str:
                        if ',' in czas_str:
                            parts = czas_str.split(',')
                            if len(parts) > 1 and parts[1] == '15': 
                                czas_num = float(parts[0]) + 0.25
                            elif len(parts) > 1 and parts[1] == '30': 
                                czas_num = float(parts[0]) + 0.5
                            elif len(parts) > 1 and parts[1] == '45': 
                                czas_num = float(parts[0]) + 0.75
                            else: 
                                czas_num = float(czas_str.replace(',', '.'))
                        else:
                            try: 
                                czas_num = float(czas_str)
                            except: 
                                czas_num = 1.0
                                
                min_start = 9.0
                max_start = 18.0
                if col_kiedy and pd.notna(match.iloc[0][col_kiedy]):
                    kiedy_str = str(match.iloc[0][col_kiedy]).strip()
                    m = re.findall(r'(\d{1,2}):\d{2}', kiedy_str)
                    if len(m) >= 2: 
                        min_start = float(m[0])
                        max_start = float(m[1])
                    elif len(m) == 1: 
                        min_start = float(m[0])
                        max_start = 22.0
                        
                return {"cena": cena, "czas": czas_num, "min_start": min_start, "max_start": max_start}
    except Exception: 
        pass
        
    return {"cena": 0, "czas": 0.0, "min_start": 9.0, "max_start": 18.0}

def sprawdz_dostepnosc_hotres(data_od, data_do):
    try:
        api_key = st.secrets["hotres"]["api"]
        auth_key = st.secrets["hotres"]["auth"]
    except KeyError: 
        return "Błąd kluczy API Hotres w pliku konfiguracyjnym secrets", {}
        
    try:
        url = f"https://panel.hotres.pl/api_availability?auth={auth_key}&apikey={api_key}"
        response = requests.get(url, timeout=15)
        
        if response.status_code == 200:
            szczegoly_zajetosci = {}
            liczba_nocy = max(1, (data_do - data_od).days)
            wymagane_daty = []
            for i in range(liczba_nocy):
                wymagane_daty.append((data_od + timedelta(days=i)).strftime("%Y-%m-%d"))
                
            for room_data in response.json():
                room_id = room_data.get("type_id")
                if str(room_id).isdigit():
                    nazwa_pokoju = HOTRES_ROOM_MAP.get(int(room_id), f"ID_{room_id}")
                else:
                    nazwa_pokoju = HOTRES_ROOM_MAP.get(room_id, f"ID_{room_id}")
                    
                for day_data in room_data.get("dates", []):
                    data_dnia = day_data.get("date")
                    dostepnosc = int(float(day_data.get("available", 0)))
                    if data_dnia in wymagane_daty and dostepnosc <= 0:
                        if nazwa_pokoju not in szczegoly_zajetosci: 
                            szczegoly_zajetosci[nazwa_pokoju] = []
                        szczegoly_zajetosci[nazwa_pokoju].append(data_dnia)
            return "", szczegoly_zajetosci
        else:
            return f"Hotres zwrócił kod błędu: {response.status_code}", {}
    except Exception as e: 
        return f"Błąd połączenia z Hotres: {str(e)}", {}

def utworz_rezerwacje_hotres(data_od, data_do, wybrane_pokoje_i_domki):
    try: 
        api_key = st.secrets["hotres"]["api"]
        auth_key = st.secrets["hotres"]["auth"]
    except KeyError: 
        return "Błąd kluczy autoryzacji Hotres."
    
    imiona = st.session_state.klient_imie.strip().split(" ", 1)
    if len(imiona) > 1:
        last_name = imiona[1]
    else:
        last_name = "-"
        
    pokoje_payload = []
    for rname in wybrane_pokoje_i_domki:
        tid = HOTRES_ROOM_NAME_TO_ID.get(rname)
        if tid: 
            pokoje_payload.append({
                "arrival_date": data_od.strftime("%Y-%m-%d"), 
                "departure_date": data_do.strftime("%Y-%m-%d"), 
                "type_id": tid, 
                "price": 0, 
                "amount": 0, 
                "adults": 1, 
                "child1": 0
            })

    if st.session_state.get("typ_klienta_radio", "Indywidualny") == "Biznesowy":
        rate_id = 31803
    else:
        rate_id = 24548

    payload = {
        "status": "new", 
        "currency": "PLN", 
        "rate_id": rate_id,
        "lang": "pl", 
        "source": "reception", 
        "first_name": imiona[0], 
        "last_name": last_name,
        "phone": st.session_state.telefon_n, 
        "phone_prefix": "48", 
        "email": st.session_state.email_n,
        "company_name": st.session_state.firma_n, 
        "company_nip": st.session_state.nip_n, 
        "rooms": pokoje_payload
    }
    
    try:
        url = f"https://panel.hotres.pl/api_reservation?auth={auth_key}&apikey={api_key}"
        res = requests.post(url, json=payload, timeout=15)
        if res.status_code in [200, 201]: 
            if res.json().get("result") == "success":
                return "OK"
            else:
                return f"Hotres odrzucił dane: {res.text}"
        else:
            return f"Błąd protokołu HTTP: {res.status_code}."
    except Exception as e: 
        return f"Błąd sieciowy API Hotres: {str(e)}"

# --- DYSK I PLIKI ---
wszystkie_pliki = []
cennik_file = None

try:
    wszystkie_pliki = fetch_all_debogora_files(ROOT_FOLDER_ID)
    if wszystkie_pliki:
        pobrano_nowe = False
        czcionki = ['Lora-Regular.ttf', 'Lora-Bold.ttf', 'PTSans-Regular.ttf', 'PTSans-Bold.ttf']
        for f in wszystkie_pliki:
            if f['name'] in czcionki and not os.path.exists(f['name']):
                try:
                    with open(f['name'], 'wb') as out: 
                        out.write(download_file(f['id']).getvalue())
                    pobrano_nowe = True
                except Exception: 
                    pass
                    
        if pobrano_nowe: 
            register_custom_fonts()
        install_fonts_for_libreoffice()
        
        c_files = []
        for f in wszystkie_pliki:
            if 'cennik' in f['name'].lower() and ('xlsx' in f['name'].lower() or 'csv' in f['name'].lower()):
                c_files.append(f)
                
        if c_files: 
            c_files_sorted = sorted(c_files, key=lambda f: 'xlsx' in f['name'].lower(), reverse=True)
            cennik_file = c_files_sorted[0]
except Exception: 
    pass

if cennik_file and 'df_cennik' not in st.session_state:
    try:
        fs = download_file(cennik_file['id'])
        if 'xlsx' in cennik_file['name'].lower():
            st.session_state.df_cennik = pd.read_excel(fs, engine='openpyxl')
        else:
            st.session_state.df_cennik = pd.read_csv(fs, encoding='utf-8')
    except Exception:
        try: 
            fs.seek(0)
            st.session_state.df_cennik = pd.read_csv(fs, encoding='cp1250')
        except: 
            st.session_state.df_cennik = None
elif 'df_cennik' not in st.session_state: 
    st.session_state.df_cennik = None

df_c = st.session_state.df_cennik

# --- MENU BOCZNE ---
with st.sidebar:
    st.header("🗂️ Karty Produktów")
    with st.expander("Rozwiń pliki do pobrania"):
        if wszystkie_pliki:
            pliki_do_pobrania = []
            for f in wszystkie_pliki:
                if 'cennik' not in f['name'].lower() and f['mimeType'] != 'application/vnd.google-apps.folder' and '.ttf' not in f['name'].lower():
                    pliki_do_pobrania.append(f)
                    
            for f in sorted(pliki_do_pobrania, key=lambda x: x['name']):
                st.markdown(f"📄 [{f['name']}]({f.get('webViewLink', '#')})")
        else: 
            st.info("Brak plików na Dysku Google.")

CENNIK = {
    "nocleg_1_noc": get_price_data("Nocleg (1 noc)", df_c)["cena"], 
    "nocleg_2_noce": get_price_data("Nocleg (2+ noce)", df_c)["cena"], 
    "doplata_domek": 40,
    "domki": {
        "Muuu 1": {"baza": get_price_data("Muuu 1, 2", df_c)["cena"], "pdf": "krovacja"},
        "Muuu 2": {"baza": get_price_data("Muuu 1, 2", df_c)["cena"], "pdf": "krovacja"},
        "Muuu 3": {"baza": get_price_data("Muuu 3, 4", df_c)["cena"], "pdf": "krovacja"},
        "Muuu 4": {"baza": get_price_data("Muuu 3, 4", df_c)["cena"], "pdf": "krovacja"},
        "Muuu 5": {"baza": get_price_data("Muuu 5, 6", df_c)["cena"], "pdf": "krovacja"},
        "Muuu 6": {"baza": get_price_data("Muuu 5, 6", df_c)["cena"], "pdf": "krovacja"}
    },
    "wyzywienie": {
        "Śniadanie": {"dane": get_price_data("Śniadanie", df_c), "pdf": "wyżywienie"},
        "Obiadokolacja": {"dane": get_price_data("Obiadokolacja", df_c), "pdf": "wyżywienie"},
        "Serwis kawowy": {"dane": get_price_data("Serwis kawowy", df_c), "pdf": "serwis kawowy"},
        "Wiejskie jadło": {"dane": get_price_data("Wiejskie jadło (Podstawowe)", df_c), "pdf": "wiejskie jadło"},
        "Kolacja z rozszerzonym menu": {"dane": get_price_data("Biesiada wieczorna", df_c), "pdf": "rozszerzonym menu"}
    },
    "SPAstwisko": {
        "Seans full experience": {"dane": get_price_data("Seans full experience", df_c), "typ": "grupa", "pdf": "Seans saunowy"},
        "Sauna olchowa": {"dane": get_price_data("Sauna Olchowa", df_c), "typ": "grupa", "pdf": "Sauna olchowa"},
        "Staw kąpielowy": {"dane": get_price_data("Staw kąpielowy", df_c), "typ": "grupa", "pdf": "Staw"},
        "Balia opalana drewnem": {"dane": get_price_data("Balia opalana drewnem", df_c), "typ": "grupa", "pdf": "Balia"},
        "Sauny": {"dane": get_price_data("Wynajem na wyłączność", df_c), "typ": "grupa", "pdf": "Sauny"},
        "Masaż relaksacyjny": {"dane": get_price_data("Masaż relaksacyjny", df_c), "typ": "osoba", "pdf": "Masaże"},
        "Masaż gorącą świecą": {"dane": get_price_data("Masaż gorącą świecą", df_c), "typ": "osoba", "pdf": "Masaże"},
    },
    "Atrakcje": {
        "Złodziej Krów": {"dane": get_price_data("Złodziej Krów", df_c), "typ": "osoba", "pdf": "Łowcy krów"},
        "Skarby Dębogóry": {"dane": get_price_data("Skarby Dębogóry", df_c), "typ": "osoba", "pdf": "Skarby"},
        "Krowie Safari Standard": {"dane": get_price_data("Krowie Safari Standard", df_c), "typ": "osoba", "pdf": "Safari_Standard"},
        "Krowie Safari Rozszerzone": {"dane": get_price_data("Krowie Safari Rozszerzone", df_c), "typ": "osoba", "pdf": "Safari_Rozszerzona"},
        "Paintball": {"dane": get_price_data("Paintball", df_c), "typ": "osoba", "pdf": "Paintball"},
        "Kajaki": {"dane": get_price_data("Kajaki", df_c), "typ": "osoba", "pdf": "Spływ kajakowy"},
        "Ognisko": {"dane": get_price_data("Ognisko", df_c), "typ": "grupa", "pdf": "Ognisko"},
        "Rowery elektryczne 1 dzień": {"dane": get_price_data("Rowery elektryczne 1 dzień", df_c), "typ": "osoba", "pdf": "Rowery"},
        "Punkt widokowy": {"dane": get_price_data("Punkt widokowy", df_c), "typ": "grupa", "pdf": "Punkt widokowy"},
    },
    "Biznes": {
        "Blok konferencyjny": {"dane": get_price_data("Blok konferencyjny", df_c), "typ": "grupa", "pdf": "Blok konferencyjny"},
        "Wynajem sali": {"dane": get_price_data("Wynajem sali", df_c), "typ": "grupa", "pdf": "Wynajem sali"},
    }
}

POKOJE_DWOREK = {}
for i in range(1, 13):
    if i == 1:
        pojemnosc = 1
    elif i == 11:
        pojemnosc = 4
    elif i in [7, 9, 10, 12]:
        pojemnosc = 3
    else:
        pojemnosc = 2
    POKOJE_DWOREK[f"Pokój nr {i}"] = pojemnosc

try:
    with open("logo.png", "rb") as image_file:
        encoded_string = base64.b64encode(image_file.read()).decode()
    st.markdown(f'<div style="display: flex; justify-content: center; margin-bottom: 10px;"><img src="data:image/png;base64,{encoded_string}" width="120"></div>', unsafe_allow_html=True)
except Exception: 
    pass

st.markdown("<h1 style='text-align: center; margin-top:0;'>PLANER OFERT</h1>", unsafe_allow_html=True)

def auto_alloc():
    total = st.session_state.l_osob_total
    zajete = list(st.session_state.szczegoly_zajetosci.keys())
    st.session_state.wybrane_p = []
    st.session_state.wybrane_d = []
            
    if st.session_state.get("typ_klienta_radio", "Indywidualny") == "Biznesowy":
        max_os_domki = {"Muuu 1": 2, "Muuu 2": 2, "Muuu 3": 4, "Muuu 4": 4, "Muuu 5": 1, "Muuu 6": 1}
    else:
        max_os_domki = {"Muuu 1": 4, "Muuu 2": 4, "Muuu 3": 6, "Muuu 4": 6, "Muuu 5": 3, "Muuu 6": 3}
        
    def p_dworek(osoby):
        for p, cap in POKOJE_DWOREK.items():
            if osoby > 0 and p not in zajete:
                st.session_state.wybrane_p.append(p)
                przydzielono = min(cap, osoby)
                st.session_state[f"os_{p}"] = przydzielono
                osoby -= przydzielono
        return osoby
        
    def p_domki(osoby):
        for d, cap in max_os_domki.items():
            if osoby > 0 and d not in zajete:
                st.session_state.wybrane_d.append(d)
                przydzielono = min(cap, osoby)
                st.session_state[f"os_{d}"] = przydzielono
                osoby -= przydzielono
        return osoby

    if st.session_state.get("marka_oferty_select", "Dwór Dębogóra") == "Dwór Dębogóra": 
        pozostale = p_dworek(total)
        p_domki(pozostale)
    else: 
        pozostale = p_domki(total)
        p_dworek(pozostale)

tab1, tab3, tab2 = st.tabs(["📝 Kreator Ofert", "📂 Baza Ofert", "⚙️ Edycja Cennika"])

with tab2:
    st.subheader("Edycja pliku Cennika (Google Drive)")
    if df_c is not None:
        edited_df = st.data_editor(df_c, num_rows="dynamic", use_container_width=True)
        if st.button("💾 ZAPISZ ZMIANY NA DYSKU", type="primary"):
            update_file_on_drive(cennik_file['id'], edited_df, cennik_file['name'])
            st.session_state.df_cennik = edited_df
            st.success("Zmiany zapisane!")

with tab3:
    st.subheader("Baza i Historia Wygenerowanych Ofert")
    if st.button("🔄 Odśwież listę archiwalną"): 
        st.rerun()
        
    try:
        if BAZA_OFERT_SHEET_ID == "TUTAJ_WKLEJ_ID_ARKUSZA": 
            st.info("⚠️ Skonfiguruj BAZA_OFERT_SHEET_ID")
        else:
            wiersze_ofert = get_offers_from_sheet(BAZA_OFERT_SHEET_ID)
            poprawne_oferty = []
            for row in wiersze_ofert:
                dopelniony_wiersz = row + [''] * (11 - len(row))
                if str(dopelniony_wiersz[0]).startswith("202"):
                    poprawne_oferty.append(dopelniony_wiersz)
                    
            if not poprawne_oferty: 
                st.info("Brak zapisanych ofert w Arkuszu Google.")
            else:
                for idx, row in enumerate(reversed(poprawne_oferty)): 
                    c_file, c_act = st.columns([3, 2])
                    
                    data = row[0]
                    imie = row[1]
                    firma = row[2]
                    osoby = row[8]
                    
                    if firma:
                        nazwa = firma
                    else:
                        nazwa = imie
                        
                    c_file.markdown(f"**📄 {data} - {nazwa} ({osoby} os.)**")
                    if c_act.button("📂 Przywróć do Kreatora", key=f"load_sheet_{idx}"):
                        st.session_state.klient_imie = row[1]
                        st.session_state.firma_n = row[2]
                        st.session_state.nip_n = row[3]
                        st.session_state.telefon_n = row[4]
                        st.session_state.email_n = row[5]
                        st.session_state.marka_oferty_select = row[6]
                        st.session_state.typ_klienta_radio = row[7]
                        
                        if str(row[8]).isdigit():
                            st.session_state.l_osob_total = int(row[8])
                        else:
                            st.session_state.l_osob_total = 10
                            
                        st.session_state.agenda_custom_text = row[9]
                        
                        # --- KLUCZOWA POPRAWKA ARCHIWUM: Odtwarzanie UI z JSONa ---
                        try: 
                            pozycje_json = json.loads(row[10])
                            st.session_state.loaded_pozycje = pozycje_json
                            
                            # Odtworzenie stanu pól wyboru (multiselect)
                            st.session_state.wybrane_p = [p["Opis"] for p in pozycje_json if p["Opis"] in POKOJE_DWOREK.keys()]
                            st.session_state.wybrane_d = [p["Opis"] for p in pozycje_json if p["Opis"] in CENNIK["domki"].keys()]
                            st.session_state.wyz_sel_key = [p["Opis"] for p in pozycje_json if p["Opis"] in CENNIK["wyzywienie"].keys()]
                            st.session_state.spa_sel_key = [p["Opis"] for p in pozycje_json if p["Opis"] in CENNIK["SPAstwisko"].keys()]
                            st.session_state.atr_sel_key = [p["Opis"] for p in pozycje_json if p["Opis"] in CENNIK["Atrakcje"].keys()]
                            st.session_state.biz_sel_key = [p["Opis"] for p in pozycje_json if p["Opis"] in CENNIK["Biznes"].keys()]
                            
                            # Odtworzenie wartości numerycznych (ilości osób/usług)
                            for p in pozycje_json:
                                o = p["Opis"]
                                i = p["Ilość"]
                                if o in POKOJE_DWOREK: st.session_state[f"os_{o}"] = i
                                elif o in CENNIK["domki"]: st.session_state[f"os_{o}"] = i
                                elif o in CENNIK["SPAstwisko"]: st.session_state[f"spa_{o}"] = i
                                elif o in CENNIK["Atrakcje"]: st.session_state[f"atr_{o}"] = i
                                elif o in CENNIK["Biznes"]: st.session_state[f"biz_{o}"] = i
                        except: 
                            st.session_state.loaded_pozycje = []
                            
                        st.rerun()
    except Exception as e: 
        st.error(f"Nie można załadować bazy: {e}")

with tab1:
    pozycje_kosztowe = []
    wybrane_atrakcje_agenda = []

    with st.container():
        c_head1, c_head2 = st.columns([4, 1])
        c_head1.subheader("1. Główne Ustawienia i Dane Klienta")
        if c_head2.button("🧹 Resetuj formularz"):
            for key in list(st.session_state.keys()):
                if key not in ['df_cennik']: 
                    del st.session_state[key]
            st.rerun()

        c1, c2 = st.columns(2)
        with c1:
            marka_oferty = st.selectbox("Marka wiodąca oferty *", ["Dwór Dębogóra", "Krovacja"], key="marka_oferty_select")
            typ_klienta = st.radio("Typ klienta", ["Indywidualny", "Biznesowy"], horizontal=True, key="typ_klienta_radio")
            st.session_state.klient_imie = st.text_input("Imię i nazwisko osoby kontaktowej *", value=st.session_state.klient_imie)
            st.session_state.firma_n = st.text_input("Firma (opcjonalnie)", value=st.session_state.firma_n)
            st.session_state.nip_n = st.text_input("NIP firmy (do integracji Hotres)", value=st.session_state.nip_n)
            st.number_input("Liczba osób", 1, 100, key="l_osob_total")
            st.button("🤖 Automatycznie rozmieść gości", on_click=auto_alloc)
            
        with c2:
            st.session_state.email_n = st.text_input("Email", value=st.session_state.email_n)
            st.session_state.telefon_n = st.text_input("Telefon kontaktowy (do integracji Hotres)", value=st.session_state.telefon_n)
            cd1, cd2 = st.columns(2)
            d_in = cd1.date_input("Przyjazd", date.today())
            d_out = cd2.date_input("Wyjazd", date.today() + timedelta(1))
            dni = max(1, (d_out - d_in).days)
            st.markdown("---")
            if st.button("🔍 Sprawdź dostępność na żywo (Hotres)"):
                with st.spinner("Łączenie z bazą rezerwacji..."):
                    err, st.session_state.szczegoly_zajetosci = sprawdz_dostepnosc_hotres(d_in, d_out)
                    if err: 
                        st.error(err)
                    elif st.session_state.szczegoly_zajetosci: 
                        st.warning(f"🚨 Wykryto zablokowane obiekty!")
                    else: 
                        st.success("✅ Wszystkie zmapowane obiekty są wolne!")

    with st.container():
        st.subheader("2. Zakwaterowanie")
        if dni == 1:
            stawka_dw = CENNIK["nocleg_1_noc"]
        else:
            stawka_dw = CENNIK["nocleg_2_noce"]
            
        col_dw, col_dm = st.columns(2)
        osoby_zadeklarowane = 0
        zajete_nazwy = list(st.session_state.szczegoly_zajetosci.keys())
        
        if typ_klienta == "Biznesowy":
            max_os_domki = {"Muuu 1": 2, "Muuu 2": 2, "Muuu 3": 4, "Muuu 4": 4, "Muuu 5": 1, "Muuu 6": 1} 
        else:
            max_os_domki = {"Muuu 1": 4, "Muuu 2": 4, "Muuu 3": 6, "Muuu 4": 6, "Muuu 5": 3, "Muuu 6": 3}
            
        with col_dw:
            pokoje_dostepne = []
            for p in POKOJE_DWOREK.keys():
                if p not in zajete_nazwy:
                    pokoje_dostepne.append(p)
                    
            p_sel = st.multiselect("Dworek (Dostępne pokoje)", pokoje_dostepne, key="wybrane_p")
            
            for p in p_sel:
                ile = st.number_input(f"{p} (Max: {POKOJE_DWOREK[p]})", 1, POKOJE_DWOREK[p], key=f"os_{p}")
                osoby_zadeklarowane += ile
                pozycje_kosztowe.append({
                    "Kategoria": "Nocleg", 
                    "Opis": f"{p}", 
                    "Ilość": ile, 
                    "Cena jednostkowa": stawka_dw * dni, 
                    "Suma": ile * stawka_dw * dni
                })
                
        with col_dm:
            domki_dostepne = []
            for d in CENNIK["domki"].keys():
                if d not in zajete_nazwy:
                    domki_dostepne.append(d)
                    
            d_sel = st.multiselect("Domki Krovacja (Dostępne)", domki_dostepne, key="wybrane_d")
            
            for d in d_sel:
                ile = st.number_input(f"{d} (Max: {max_os_domki[d]})", 1, max_os_domki[d], key=f"os_{d}")
                osoby_zadeklarowane += ile
                cena_d = (CENNIK["domki"][d]["baza"] + (max(0, ile - 1) * CENNIK["doplata_domek"])) * dni
                pozycje_kosztowe.append({
                    "Kategoria": "Nocleg", 
                    "Opis": f"{d}", 
                    "Ilość": 1, 
                    "Cena jednostkowa": cena_d, 
                    "Suma": cena_d
                })

        if osoby_zadeklarowane > st.session_state.l_osob_total:
            overbooking_error = True
            st.error(f"⚠️ Przydzieliłeś {osoby_zadeklarowane} miejsc dla {st.session_state.l_osob_total} gości.")
        else:
            overbooking_error = False

    with st.container():
        st.subheader("3. Wyżywienie")
        wyz_sel = st.multiselect("Wybierz opcje wyżywienia", list(CENNIK["wyzywienie"].keys()), key="wyz_sel_key")
        for w in wyz_sel:
            ile = st.number_input(f"Ilość porcji: {w}", 1, 5000, st.session_state.l_osob_total * dni)
            cena_w = CENNIK["wyzywienie"][w]["dane"]["cena"]
            pozycje_kosztowe.append({
                "Kategoria": "Gastronomia", 
                "Opis": w, 
                "Ilość": ile, 
                "Cena jednostkowa": cena_w, 
                "Suma": ile * cena_w
            })

    with st.container():
        st.subheader("4. Oferta Dodatkowa (SPAstwisko, Atrakcje, Biznes)")
        c_spa, c_atr, c_biz = st.columns(3)
        
        def render_dodatki(kolumna, tytul, cennik_klucz, prefix, session_key):
            opcje = list(CENNIK[cennik_klucz].keys())
            sel = kolumna.multiselect(tytul, opcje, key=session_key)
            for a in sel:
                dane = CENNIK[cennik_klucz][a]
                if dane["typ"] == "osoba":
                    def_ilo = st.session_state.l_osob_total
                else:
                    def_ilo = 1
                    
                ile = kolumna.number_input(f"Ilość: {a}", 1, 100, def_ilo, key=f"{prefix}_{a}")
                
                cena = dane["dane"]["cena"]
                pozycje_kosztowe.append({
                    "Kategoria": tytul, 
                    "Opis": a, 
                    "Ilość": ile, 
                    "Cena jednostkowa": cena, 
                    "Suma": ile * cena
                })
                wybrane_atrakcje_agenda.append({
                    "Nazwa": a, 
                    "Czas": dane["dane"]["czas"], 
                    "MinStart": dane["dane"]["min_start"], 
                    "MaxStart": dane["dane"]["max_start"]
                })
            return sel
            
        spa_sel = render_dodatki(c_spa, "SPAstwisko", "SPAstwisko", "spa", "spa_sel_key")
        atr_sel = render_dodatki(c_atr, "Atrakcje", "Atrakcje", "atr", "atr_sel_key")
        biz_sel = render_dodatki(c_biz, "Biznes", "Biznes", "biz", "biz_sel_key")

    with st.container():
        st.subheader("5. Generator Harmonogramu (Agenda)")
        roznica_dni = (d_out - d_in).days
        if roznica_dni > 0:
            liczba_dni = roznica_dni + 1
        else:
            liczba_dni = 1
        
        def format_time(hours_float):
            h = int(hours_float)
            m = int(round((hours_float - h) * 60))
            return f"{h:02d}:{m:02d}"

        unassigned = wybrane_atrakcje_agenda.copy()
        unassigned.sort(key=lambda x: (x['MaxStart'] - x['MinStart'], x['MinStart']))
        
        def fill_slot(slot_start_h, slot_end_h, pending_events):
            events = []
            curr_h = slot_start_h
            while pending_events and curr_h < slot_end_h:
                best_idx, best_start, best_end = -1, -1, -1
                for i, atr in enumerate(pending_events):
                    prop_start = max(curr_h, atr['MinStart'])
                    if atr['Czas'] > 0:
                        czas = atr['Czas']
                    else:
                        czas = 1.0
                    prop_end = prop_start + czas
                    
                    if prop_start <= atr['MaxStart'] and prop_end <= slot_end_h:
                        best_idx = i
                        best_start = prop_start
                        best_end = prop_end
                        break
                        
                if best_idx != -1:
                    atr = pending_events.pop(best_idx)
                    if best_start > curr_h: 
                        events.append({"Nazwa": "Czas wolny", "Start": curr_h, "End": best_start})
                    events.append({"Nazwa": atr["Nazwa"], "Start": best_start, "End": best_end})
                    curr_h = best_end
                else: 
                    break 
            return events, pending_events
            
        def render_events(events): 
            txt = ""
            for e in events:
                txt += f"• {format_time(e['Start'])} - {format_time(e['End'])} : {e['Nazwa']}\n"
            return txt

        nowy_draft_agendy = f"TERMIN WYDARZENIA: {d_in.strftime('%d.%m.%Y')} - {d_out.strftime('%d.%m.%Y')}\n\n"
        
        # --- KLUCZOWA POPRAWKA AGENDY: Sprawdzanie wyboru posiłków ---
        ma_sniadanie = "Śniadanie" in wyz_sel
        ma_obiad = "Obiadokolacja" in wyz_sel

        if liczba_dni == 1:
            evs, unassigned = fill_slot(10.5, 18.0, unassigned)
            evs_eve, unassigned = fill_slot(19.0, 23.0, unassigned)
            nowy_draft_agendy += f"DZIEŃ 1 ({d_in.strftime('%d.%m')})\n"
            nowy_draft_agendy += f"• 10:00 - Przyjazd\n"
            nowy_draft_agendy += render_events(evs)
            if ma_obiad:
                nowy_draft_agendy += f"• 18:00 - 19:00 : Obiadokolacja\n"
            nowy_draft_agendy += render_events(evs_eve)
            nowy_draft_agendy += f"• 23:00 - Zakończenie pobytu\n"
        else:
            evs, unassigned = fill_slot(16.0, 18.0, unassigned) 
            evs_eve, unassigned = fill_slot(19.0, 23.0, unassigned) 
            nowy_draft_agendy += f"DZIEŃ 1 ({d_in.strftime('%d.%m')})\n"
            nowy_draft_agendy += f"• 15:00 - Przyjazd i Zakwaterowanie\n"
            nowy_draft_agendy += render_events(evs)
            if ma_obiad:
                nowy_draft_agendy += f"• 18:00 - 19:00 : Obiadokolacja\n"
            nowy_draft_agendy += render_events(evs_eve)
            nowy_draft_agendy += "\n"
            
            for d in range(2, liczba_dni):
                evs, unassigned = fill_slot(10.0, 18.0, unassigned) 
                evs_eve, unassigned = fill_slot(19.0, 23.0, unassigned)
                data_dnia = (d_in + timedelta(days=d-1)).strftime('%d.%m')
                nowy_draft_agendy += f"DZIEŃ {d} ({data_dnia})\n"
                if ma_sniadanie:
                    nowy_draft_agendy += f"• 09:00 - 10:00 : Śniadanie\n"
                nowy_draft_agendy += render_events(evs)
                if ma_obiad:
                    nowy_draft_agendy += f"• 18:00 - 19:00 : Obiadokolacja\n"
                nowy_draft_agendy += render_events(evs_eve)
                nowy_draft_agendy += "\n"
                
            evs, unassigned = fill_slot(10.0, 13.0, unassigned) 
            nowy_draft_agendy += f"DZIEŃ {liczba_dni} ({d_out.strftime('%d.%m')}) (Wyjazd)\n"
            if ma_sniadanie:
                nowy_draft_agendy += f"• 09:00 - 10:00 : Śniadanie\n"
            nowy_draft_agendy += render_events(evs)
            nowy_draft_agendy += f"• 13:00 - Wymeldowanie\n"

        # Dodajemy wyz_sel do skrótu hash, żeby po usunięciu śniadania, okienko wygenerowało się na nowo
        aktualne_parametry_agendy = f"{d_in}_{d_out}_{wybrane_atrakcje_agenda}_{wyz_sel}"
        if st.session_state.get("parametry_agendy_hash") != aktualne_parametry_agendy:
            st.session_state.agenda_custom_text = nowy_draft_agendy
            st.session_state.parametry_agendy_hash = aktualne_parametry_agendy

        st.session_state.agenda_custom_text = st.text_area("Szkic Harmonogramu (do edycji):", value=st.session_state.agenda_custom_text, height=350)

    # --- ZARZĄDZANIE STRONAMI PDF ---
    def generate_default_pdf_tags():
        tags = ["okładka", "powitalna"]
        has_dworek = False
        if st.session_state.wybrane_p:
            has_dworek = True
            
        has_krovacja = False
        if st.session_state.wybrane_d:
            has_krovacja = True
        
        if has_dworek and has_krovacja: 
            tags.extend(["zakwaterowanie_oba", "uklad_debogora", "uklad_krovacja"])
        elif marka_oferty == "Krovacja" or has_krovacja: 
            tags.extend(["zakwaterowanie_domki", "uklad_krovacja"])
        else: 
            tags.extend(["zakwaterowanie_dwor", "uklad_debogora"])

        for w in wyz_sel: 
            tags.append(CENNIK["wyzywienie"][w]["pdf"])
            
        tags.append("atrakcje_wstęp")
        
        for a in spa_sel: 
            tags.append(CENNIK["SPAstwisko"][a]["pdf"])
        for a in atr_sel: 
            tags.append(CENNIK["Atrakcje"][a]["pdf"])
        for a in biz_sel: 
            tags.append(CENNIK["Biznes"][a]["pdf"])
        
        tags.extend(["wycena", "agenda", "kontakt"])
        
        seen = set()
        final_tags = []
        for x in tags:
            if x not in seen:
                seen.add(x)
                final_tags.append(x)
        return final_tags

    aktualne_parametry_ofert = f"{marka_oferty}_{st.session_state.wybrane_p}_{st.session_state.wybrane_d}_{wyz_sel}_{spa_sel}_{atr_sel}_{biz_sel}"
    if st.session_state.get("parametry_ofert_hash") != aktualne_parametry_ofert:
        st.session_state.pdf_page_list = generate_default_pdf_tags()
        st.session_state.parametry_ofert_hash = aktualne_parametry_ofert

    with st.container():
        st.subheader("6. Kosztorys, Rezerwacja Hotres i Eksport PDF")
        
        # --- TABELA KOSZTOWA ---
        ui_hash = hash(str(pozycje_kosztowe))
        if st.session_state.loaded_pozycje is not None:
            st.session_state.aktualna_tabela = pd.DataFrame(st.session_state.loaded_pozycje)
            st.session_state.loaded_pozycje = None
            st.session_state.ostatni_ui_hash = ui_hash
        elif "ostatni_ui_hash" not in st.session_state or ui_hash != st.session_state.ostatni_ui_hash:
            st.session_state.aktualna_tabela = pd.DataFrame(pozycje_kosztowe)
            st.session_state.ostatni_ui_hash = ui_hash
            
        df_robocze = st.session_state.aktualna_tabela.copy()
        if not df_robocze.empty:
            df_robocze = df_robocze[["Kategoria", "Opis", "Ilość", "Cena jednostkowa", "Suma"]]
            st.info("💡 Zmieniłeś cenę lub ilość? Kliknij **Przelicz Tabelę**, aby zaktualizować sumy we wszystkich wierszach.")
            
            edf = st.data_editor(
                df_robocze, 
                use_container_width=True, 
                num_rows="dynamic", 
                column_config={
                    "Suma": st.column_config.NumberColumn("Suma (Wciśnij Przelicz)", disabled=True, format="%.2f PLN"),
                    "Ilość": st.column_config.NumberColumn("Ilość", min_value=0.0),
                    "Cena jednostkowa": st.column_config.NumberColumn("Cena jedn.", min_value=0.0, format="%.2f PLN")
                }
            )
            st.session_state.aktualna_tabela = edf.copy()
            
            if st.button("🧮 PRZELICZ TABELĘ Z KOSZTORYSEM", type="secondary", use_container_width=True):
                edf["Ilość"] = pd.to_numeric(edf["Ilość"], errors='coerce').fillna(0)
                edf["Cena jednostkowa"] = pd.to_numeric(edf["Cena jednostkowa"], errors='coerce').fillna(0)
                edf["Suma"] = edf["Ilość"] * edf["Cena jednostkowa"]
                st.session_state.aktualna_tabela = edf.copy()
                st.rerun()
            
            edf_pdf = edf.copy()
            edf_pdf["Ilość"] = pd.to_numeric(edf_pdf["Ilość"], errors='coerce').fillna(0)
            edf_pdf["Cena jednostkowa"] = pd.to_numeric(edf_pdf["Cena jednostkowa"], errors='coerce').fillna(0)
            edf_pdf["Suma"] = edf_pdf["Ilość"] * edf_pdf["Cena jednostkowa"]
            razem = edf_pdf["Suma"].sum()
            
            st.markdown(f"<h3 style='color: {CI['dark_green']}; text-align: center; margin-top: 15px;'>RAZEM DO ZAPŁATY: {razem:,.2f} PLN</h3>", unsafe_allow_html=True)
            st.markdown("---")

            # --- MENEDŻER STRON PDF ---
            st.markdown("#### Menedżer Stron PDF")
            st.info("Oto karty, które zostaną skompilowane do finalnego pliku PDF. Kolejność na liście to kolejność stron w dokumencie. **Możesz dodawać własne strony, usuwać je lub zamieniać kolejność (edytując numery).**")
            
            lista_stron_do_tabeli = []
            for i, strona in enumerate(st.session_state.pdf_page_list):
                lista_stron_do_tabeli.append({"Kolejność": i + 1, "Strona z Dysku": strona})
                
            df_pages = pd.DataFrame(lista_stron_do_tabeli)
            
            edited_pages = st.data_editor(
                df_pages, 
                num_rows="dynamic", 
                use_container_width=True,
                column_config={
                    "Kolejność": st.column_config.NumberColumn("Numer Strony", min_value=1, step=1),
                    "Strona z Dysku": st.column_config.SelectboxColumn("Wybierz kartę", options=list(SYNONYMS.keys()), required=True)
                }
            )
            
            zatwierdzone_strony_pdf = edited_pages.sort_values("Kolejność")["Strona z Dysku"].dropna().tolist()

            c_actions1, c_actions2 = st.columns(2)
            with c_actions1:
                if st.button("GENERUJ FINALNY PDF Oferty", disabled=overbooking_error, type="primary"):
                    if not st.session_state.klient_imie: 
                        st.error("Podaj imię i nazwisko klienta!")
                    else:
                        with st.spinner("Kompilowanie oferty (izolacja sesji aktywowana)..."):
                            session_uid = uuid.uuid4().hex
                            try:
                                merger = PdfWriter()
                                open_streams = []
                                missing_cards = []
                                added_file_ids = set()
                                
                                has_dworek = False
                                if st.session_state.wybrane_p:
                                    has_dworek = True
                                    
                                has_krovacja = False
                                if st.session_state.wybrane_d:
                                    has_krovacja = True
                                    
                                if has_dworek and has_krovacja:
                                    zakwaterowanie_txt = "domkach i pokojach"
                                elif marka_oferty == "Krovacja":
                                    zakwaterowanie_txt = "domkach"
                                else:
                                    zakwaterowanie_txt = "pokojach"
                                    
                                atr_list = []
                                for _, row in edf_pdf.iterrows():
                                    if row["Kategoria"] in ["SPAstwisko", "Atrakcje", "Biznes"]:
                                        atr_list.append(row["Opis"])
                                        
                                if len(atr_list) >= 2:
                                    atrakcje_txt = f"{atr_list[0]} oraz {atr_list[1]}"
                                elif len(atr_list) == 1:
                                    atrakcje_txt = f"{atr_list[0]} oraz naturę"
                                else:
                                    atrakcje_txt = "spokój i bliskość natury"
                                
                                if st.session_state.firma_n:
                                    docelowa_nazwa = st.session_state.firma_n
                                else:
                                    docelowa_nazwa = st.session_state.klient_imie
                                    
                                if marka_oferty == "Krovacja":
                                    nazwa_obiektu = "Krovację"
                                else:
                                    nazwa_obiektu = "Dwór Dębogóra"
                                    
                                replacements = {
                                    "{{nazwa firmy}}": docelowa_nazwa, 
                                    "{{Dwór Dębogóra/Krovację}}": nazwa_obiektu, 
                                    "{{domkach/pokojach}}": zakwaterowanie_txt,
                                    "{{atrakcja}} oraz {{atrakcja}}": atrakcje_txt, 
                                    "{{przykład agendy}}": st.session_state.agenda_custom_text, 
                                    "{{tabela z wyceną dla b2b }}": "", 
                                    "{{ tabela z wyceną dla b2b }}": ""
                                }
                                
                                for tag_strony in zatwierdzone_strony_pdf:
                                    if tag_strony == "wycena":
                                        buf = io.BytesIO()
                                        doc = SimpleDocTemplate(buf, pagesize=A4, rightMargin=40, leftMargin=40, topMargin=120, bottomMargin=50)
                                        t_data = [["Kategoria", "Opis usługi", "Ilość", "Suma"]]
                                        
                                        for _, row in edf_pdf.iterrows(): 
                                            t_data.append([
                                                safe_str(row["Kategoria"]), 
                                                safe_str(row["Opis"]), 
                                                safe_str(row["Ilość"]), 
                                                f"{row['Suma']:,.0f} zł".replace(",", " ")
                                            ])
                                            
                                        t_data.append(["", "", "RAZEM:", f"{razem:,.0f} zł".replace(",", " ")])
                                        
                                        # Dodano `repeatRows=1` aby ładnie łamało strony tabeli
                                        table = Table(t_data, colWidths=[110, 220, 50, 90], repeatRows=1)
                                        table.setStyle(TableStyle([
                                            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(CI['dark_green'])), 
                                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                                            ('FONTNAME', (0, 0), (-1, 0), FONT_HEADER), 
                                            ('FONTSIZE', (0, 0), (-1, 0), 12), 
                                            ('ALIGN', (0, 0), (-1, -1), 'CENTER'), 
                                            ('ALIGN', (1, 1), (1, -2), 'LEFT'), 
                                            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'), 
                                            ('BOTTOMPADDING', (0, 0), (-1, -1), 10), 
                                            ('TOPPADDING', (0, 0), (-1, -1), 10), 
                                            ('FONTNAME', (0, 1), (-1, -1), FONT_TEXT), 
                                            ('FONTNAME', (2, -1), (-1, -1), FONT_TEXT_BOLD),
                                            ('TEXTCOLOR', (2, -1), (-1, -1), colors.HexColor(CI['dark_green'])), 
                                            ('BACKGROUND', (2, -1), (-1, -1), colors.HexColor(CI['light_green'])),
                                            ('LINEABOVE', (0, -1), (-1, -1), 1, colors.HexColor(CI['dark_green'])),
                                        ]))
                                        doc.build([table])
                                        buf.seek(0)
                                        
                                        wycena_file = get_file_by_keyword("wycena", wszystkie_pliki)
                                        if wycena_file:
                                            try:
                                                temp_ppt_w = f"temp_wycena_{session_uid}.pptx"
                                                with open(temp_ppt_w, "wb") as f: 
                                                    f.write(download_file(wycena_file['id']).getvalue())
                                                    
                                                prs = Presentation(temp_ppt_w)
                                                replace_text_in_pptx(prs, replacements)
                                                prs.save(temp_ppt_w)
                                                
                                                lo_profile_dir = f"./lo_profile_{session_uid}"
                                                lo_profile_flag = f"-env:UserInstallation=file://{os.path.abspath(lo_profile_dir)}"
                                                
                                                subprocess.run(["libreoffice", lo_profile_flag, "--headless", "--convert-to", "pdf", temp_ppt_w], check=True)
                                                
                                                temp_pdf_w = f"temp_wycena_{session_uid}.pdf"
                                                with open(temp_pdf_w, "rb") as f: 
                                                    bg_bytes = f.read()
                                                
                                                fg_reader = PdfReader(buf)
                                                
                                                # --- KLUCZOWA POPRAWKA PAGINACJI TABELI PDF ---
                                                for i, fg_page in enumerate(fg_reader.pages):
                                                    # Za każdym razem ładujemy stronę tła "na czysto"
                                                    fresh_bg_reader = PdfReader(io.BytesIO(bg_bytes))
                                                    
                                                    if i < len(fresh_bg_reader.pages):
                                                        bg_page = fresh_bg_reader.pages[i]
                                                    else:
                                                        bg_page = fresh_bg_reader.pages[-1]
                                                        
                                                    bg_page.merge_page(fg_page)
                                                    merger.add_page(bg_page)
                                                    
                                            except Exception: 
                                                merger.append(PdfReader(buf, strict=False))
                                        else: 
                                            merger.append(PdfReader(buf, strict=False))
                                    else:
                                        add_file_to_merger(merger, tag_strony, wszystkie_pliki, open_streams, missing_cards, added_file_ids, session_uid, replacements)

                                final_pdf = io.BytesIO()
                                merger.write(final_pdf)
                                pdf_bytes_to_upload = final_pdf.getvalue()
                                timestamp = datetime.now().strftime('%d%m%H%M')
                                nazwa_pliku_pdf = f"Oferta_{safe_str(st.session_state.klient_imie).replace(' ', '_')}_{timestamp}.pdf"
                                
                                st.success("✅ Oferta w formacie PDF została pomyślnie wygenerowana!")
                                st.download_button("📥 POBIERZ PDF NA DYSK LOKALNY", pdf_bytes_to_upload, nazwa_pliku_pdf, "application/pdf", type="primary")

                                pozycje_do_arkusza = []
                                for index, wiersz in edf_pdf.iterrows():
                                    pozycje_do_arkusza.append({
                                        "Kategoria": wiersz["Kategoria"],
                                        "Opis": wiersz["Opis"],
                                        "Ilość": wiersz["Ilość"],
                                        "Cena jednostkowa": wiersz["Cena jednostkowa"],
                                        "Suma": wiersz["Suma"]
                                    })

                                meta_payload = {
                                    "klient_imie": st.session_state.klient_imie, 
                                    "firma_n": st.session_state.firma_n, 
                                    "nip_n": st.session_state.nip_n, 
                                    "telefon_n": st.session_state.telefon_n, 
                                    "email_n": st.session_state.email_n, 
                                    "marka_oferty": marka_oferty,
                                    "typ_klienta": typ_klienta, 
                                    "l_osob_total": st.session_state.l_osob_total, 
                                    "final_agenda_text": st.session_state.agenda_custom_text, 
                                    "pozycje": pozycje_do_arkusza
                                }
                                
                                if BAZA_OFERT_SHEET_ID != "TUTAJ_WKLEJ_ID_ARKUSZA": 
                                    save_offer_to_sheet(BAZA_OFERT_SHEET_ID, meta_payload)
                                    
                                upload_file_to_drive(pdf_bytes_to_upload, nazwa_pliku_pdf, BAZA_OFERT_FOLDER_ID, 'application/pdf')
                                    
                            except Exception as e: 
                                st.error(f"❌ Błąd generatora: {str(e)}")
                            finally:
                                pliki_do_usuniecia = glob.glob(f"*{session_uid}*")
                                for f in pliki_do_usuniecia:
                                    try: 
                                        if os.path.isdir(f): 
                                            shutil.rmtree(f)
                                        else: 
                                            os.remove(f)
                                    except: 
                                        pass
            
            with c_actions2:
                if st.button("⚡ PRZEŚLIJ REZERWACJĘ DO HOTRES", type="secondary", use_container_width=True):
                    wybrane_obiekty = []
                    wybrane_obiekty.extend(st.session_state.wybrane_p)
                    wybrane_obiekty.extend(st.session_state.wybrane_d)
                    
                    if not wybrane_obiekty: 
                        st.error("Wybierz pokoje/domki!")
                    elif not st.session_state.klient_imie or not st.session_state.email_n: 
                        st.error("Wymagane imię, nazwisko oraz email!")
                    else:
                        with st.spinner("Wysyłanie do Hotres..."):
                            status = utworz_rezerwacje_hotres(d_in, d_out, wybrane_obiekty)
                            if status == "OK": 
                                st.success("🎉 Sukces! Pokoje zablokowane w Hotres.")
                            else: 
                                st.error(status)
